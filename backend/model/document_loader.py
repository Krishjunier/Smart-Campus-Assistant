import logging
from pathlib import Path
from typing import List
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader, TextLoader, UnstructuredWordDocumentLoader

logger = logging.getLogger(__name__)

# Check for optional dependencies
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not installed. Table extraction will be limited.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PowerPoint support unavailable.")


class EnhancedDocumentLoader:
    """Enhanced document loader supporting multiple educational formats"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.path = Path(file_path)
    
    def load(self) -> List[Document]:
        """Load document based on file type"""
        extension = self.path.suffix.lower()
        
        try:
            if extension == '.pdf':
                return self._load_pdf()
            elif extension in ['.txt', '.md']:
                return self._load_text()
            elif extension in ['.doc', '.docx']:
                return self._load_word()
            elif extension in ['.ppt', '.pptx']:
                return self._load_powerpoint()
            else:
                logger.warning(f"Unsupported file format: {extension}")
                return []
        except Exception as e:
            logger.error(f"Error loading {self.file_path}: {e}")
            return []
    
    def _load_pdf(self) -> List[Document]:
        """Load PDF with enhanced extraction"""
        if PDFPLUMBER_AVAILABLE:
            return self._load_pdf_with_pdfplumber()
        else:
            loader = PyPDFLoader(self.file_path)
            return loader.load()
    
    def _load_pdf_with_pdfplumber(self) -> List[Document]:
        """Load PDF using pdfplumber"""
        import pdfplumber
        
        documents = []
        
        with pdfplumber.open(self.file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                tables = page.extract_tables()
                
                # Combine text and tables
                combined_content = text
                
                if tables:
                    for i, table in enumerate(tables):
                        if table:
                            table_text = self._format_table(table, i)
                            combined_content += f"\n\n{table_text}"
                
                if combined_content.strip():
                    doc = Document(
                        page_content=combined_content,
                        metadata={
                            "source": self.file_path,
                            "source_file": self.path.name,
                            "page": page_num + 1,
                            "file_type": "pdf",
                            "has_tables": len(tables) > 0,
                        }
                    )
                    documents.append(doc)
        
        return documents
    
    def _format_table(self, table: List[List[str]], index: int) -> str:
        """Format table for better readability"""
        if not table:
            return ""
        
        lines = [f"[TABLE {index + 1}]"]
        for i, row in enumerate(table):
            cleaned_row = [str(cell or "").strip() for cell in row]
            if i == 0:
                lines.append("Headers: " + " | ".join(cleaned_row))
                lines.append("-" * 50)
            else:
                lines.append(" | ".join(cleaned_row))
        lines.append("[END TABLE]")
        
        return "\n".join(lines)
    
    def _load_text(self) -> List[Document]:
        """Load text files"""
        loader = TextLoader(self.file_path, encoding='utf-8')
        return loader.load()
    
    def _load_word(self) -> List[Document]:
        """Load Word documents"""
        loader = UnstructuredWordDocumentLoader(self.file_path)
        return loader.load()
    
    def _load_powerpoint(self) -> List[Document]:
        """Load PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not installed. Cannot load PowerPoint files.")
            return []
        
        from pptx import Presentation
        
        # Handle .ppt files (old format) - try to use unstructured library as fallback
        if self.path.suffix.lower() == '.ppt':
            try:
                from unstructured.partition.ppt import partition_ppt
                elements = partition_ppt(filename=self.file_path)
                
                documents = []
                slide_content = []
                
                for element in elements:
                    text = str(element)
                    if text.strip():
                        slide_content.append(text)
                
                # Combine all content into one document for .ppt
                if slide_content:
                    content = "\n\n".join(slide_content)
                    doc = Document(
                        page_content=content,
                        metadata={
                            "source": self.file_path,
                            "source_file": self.path.name,
                            "file_type": "ppt",
                            "note": "Converted from legacy .ppt format"
                        }
                    )
                    documents.append(doc)
                
                return documents
                
            except Exception as e:
                logger.error(f"Error loading .ppt file with unstructured: {e}")
                logger.info("Attempting alternative method for .ppt file...")
                logger.warning(f"Cannot fully process .ppt file. Please convert to .pptx format for better results.")
                return []
        
        # Handle .pptx files (new format)
        documents = []
        prs = Presentation(self.file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            content_parts = []
            
            # Extract title
            if slide.shapes.title:
                content_parts.append(f"SLIDE TITLE: {slide.shapes.title.text}")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content_parts.append(shape.text)
            
            if content_parts:
                content = "\n\n".join(content_parts)
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": self.file_path,
                        "source_file": self.path.name,
                        "slide": slide_num + 1,
                        "file_type": "pptx",
                    }
                )
                documents.append(doc)
        
        return documents
